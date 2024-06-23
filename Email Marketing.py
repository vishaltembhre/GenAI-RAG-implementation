import os,sys
import streamlit as st
from langchain.chat_models import AzureChatOpenAI
from langchain.embeddings import OpenAIEmbeddings
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.text_splitter import CharacterTextSplitter
from langchain.vectorstores import FAISS
import pandas as pd
# from langchain.agents import create_csv_agent
#from langchain_experimental.agents import create_csv_agent
from PyPDF2 import PdfReader
from pptx import Presentation
from azure.identity import ClientSecretCredential
from azure.keyvault.secrets import SecretClient
from langchain.agents import create_pandas_dataframe_agent
import matplotlib as plt
import seaborn as sns


def write_text_file(productDescFile,file_path,uploadProdtype):
    if uploadProdtype == 'TEXT':
        try:
            content = productDescFile.read().decode('utf-8')
            with open(file_path, 'w') as file:
                file.write(content)
                # st.success("text Uploaded!!")
            return True
        except Exception as e:
            print(f"Error occurred while writing the file: {e}")
            return False
    elif uploadProdtype == 'PDF':
        try:
            with open(productDescFile, "rb") as pdf_file:
                pdf_reader = PdfReader(pdf_file)
                # Extract text from all pages (modify for specific pages)
                content = ""
                for page in pdf_reader.pages:
                    content += page.extract_text()
                with open(file_path, 'w') as file:
                    file.write(content)
            return True
        except Exception as e:
            print(f"Error occurred while writing the file: {e}")
            return False
    elif uploadProdtype == 'PPT':
        try:
            presentation = Presentation(productDescFile)
            content = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text:
                        content += shape.text_frame.text
            # Write the data to a new text file (optional)
            with open(file_path, 'w') as file:
                file.write(content)
            return True
        except Exception as e:
            print(f"Error occurred while writing the file: {e}")
            return False
    elif uploadProdtype == 'CSV' or 'EXCEL':
        print('axc')
        if uploadProdtype == 'Excel':
            df_csv = pd.read_excel(productDescFile, encoding= 'unicode_escape')
        else:
            df_csv = pd.read_csv(productDescFile, encoding= 'unicode_escape')


        # df_csv_top10 = df_csv.head(2)
        # st.write("Sample records displayed...")
        # st.dataframe(df_csv_top10)
       
        file_path = "./temp/file.csv"
        df_csv.to_csv(file_path)
        # st.success("File Loaded Successfully!!")
   
def emailPrompt(file_path,emailMarketingTyp,legalFooter,userInputEmail,emailLength,creativity,insights):


    llm = AzureChatOpenAI(deployment_name="gpt-4", model_name="gpt-4")
    embeddings = OpenAIEmbeddings(deployment="embeddings", chunk_size=1)


    emaillen = {"Short": 100,"Medium": 200,"Long": 300}


    # Open the text file and read the text.
    loader = open(file_path, "r")
    docs = loader.read()
    #text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
   
    text_splitter = CharacterTextSplitter(        
    separator = "\n",
    chunk_size = 500,
    chunk_overlap  = 100,
    length_function = len,
    )
   
    texts = text_splitter.split_text(docs)
    docSearch = FAISS.from_texts(texts, embeddings)
    # st.success("File Loaded Successfully!!")


    similar_doc = docSearch.similarity_search(userInputEmail, k=1)
    context = similar_doc[0].page_content


    emailMarketingTyp = str(emailMarketingTyp)
    # context = str(context)
    legalFooter = str(legalFooter)
    emailLenn = str(emaillen[emailLength])
    userInputEmail = str(userInputEmail)
    insights = str(insights)


    # Prompt
    prompt_template = f"""Use the following pieces of context to answer the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer.


    Context(Start):
    - Type of Email Marketing: {{emailMarketingTyp}}


    - Product Description: {{context}}


    - Email Legal Outro: {{legalFooter}} [Append as it is and update only if grammar , syntax or semantics need to be updated and separate the legal outrow with email body using ----------- symbol]


    - Email Length - {{emailLenn}} Words


    - insight from input data - {{insights}}
    Context(End)


    - Creativity on the scale of 0 to 10 - {{creativity}}
   
    User Request: {{userInputEmail}}


    Question: Please generate an email for marketing based on the User request using the context?
    Answer:
    """
    prompt = PromptTemplate(template=prompt_template, input_variables=["emailMarketingTyp", "context","legalFooter","emailLenn","creativity","userInputEmail","insights"])


    query_llm = LLMChain(llm=llm, prompt=prompt)
    response = query_llm.run({"emailMarketingTyp": emailMarketingTyp,"context": context, "legalFooter": legalFooter,"emailLenn":emailLenn,"creativity":creativity,"userInputEmail": userInputEmail,"insights":insights})        
    #response = chain.run(input_documents=similar_doc[0].page_content, question=question)
    # st.write(response)
    return response

def legalPrompt(outputEmail,emailLength):
    llm = AzureChatOpenAI(deployment_name="gpt-4", model_name="gpt-4")


    emaillen = {"Short": 100,"Medium": 200,"Long": 300}
    emailLenn = str(emaillen[emailLength])
    outputEmail = str(outputEmail)


    # Prompt
    prompt_template = f"""
        Ensure the pharmaceutical marketing email complies with all relevant legal and regulatory requirements before it is shared with doctors residing in the USA.
        IMPORTANT - Output of the prompt needs to be the Marketing Email as provided in below Email Content with required correction if needed also the length of email is provided and it needs to within limits
        and incase any correction please specify the reason below the marketing email seperated by ------------- line.


        Email Content
        {{outputEmail}}


        Email Length - {{emailLenn}} Words


        Ensure that all claims about the drug are truthful, non-misleading, and supported by substantial evidence.
        Verify that any risk information (side effects, contraindications, etc.) is presented with equal prominence and readability as the benefit information.
        HIPAA Compliance


        Confirm that no patient-specific information is included in the email content or attachments.
        Ensure that the email content respects patient privacy and does not inadvertently disclose protected health information (PHI).
        CAN-SPAM Act


        Check that the subject line is not deceptive or misleading.
        Verify that the email includes a clear and conspicuous identification that it is an advertisement.
        Ensure that the email contains a valid physical postal address of the sender.
        Include a clear and easy-to-use opt-out mechanism.
        PhRMA Code


        Verify compliance with any additional state-specific regulations where the email recipients reside.
        Review Checklist
        Accuracy: All medical claims are accurate and supported by evidence.
        Balance: Risk and benefit information is balanced and presented with equal prominence.
        Privacy: No PHI or patient-specific information is disclosed.
        Transparency: The email is clearly marked as promotional, and the sender’s identity is clear.
        Ethics: Content aligns with ethical guidelines for interactions with healthcare professionals.
    """
    prompt = PromptTemplate(template=prompt_template, input_variables=["outputEmail","emailLenn"])


    query_llm = LLMChain(llm=llm, prompt=prompt)
    finalEmail = query_llm.run({"outputEmail": outputEmail,"emailLenn": emailLenn })        
    #response = chain.run(input_documents=similar_doc[0].page_content, question=question)
    # st.write(response)
    return finalEmail

def dataPrompt(file_pathData, userInputData,graph_path):
    # Load the CSV data into a pandas DataFrame
    llm = AzureChatOpenAI(deployment_name="gpt-4", model_name="gpt-4")
    df = pd.read_csv(file_pathData)
    df_data = df[:100]


    prompt_template = """
    You are an matplotlib expert. Below is given dataframe:
    {df_data}
    and user instruction : {userInputData}
    Please generate the necessary Python code to create a meaningful matplotlib graph based on this data. Ensure the code is executable in a standard Python environment with matplotlib.
    Do not add comments or ```, just output the code."""
    prompt = PromptTemplate(template=prompt_template, input_variables=["df_data","userInputData"])
    graph_llm = LLMChain(llm=llm, prompt=prompt)
    response_graph = graph_llm.run({"df_data": df_data,"userInputData" : userInputData})


    try:
        exec_context = {}
        exec(response_graph, {}, exec_context)
        graph_path = "./temp/plot.png"
        exec_context['plt'].savefig(graph_path)
        exec_context['plt'].close()
    except Exception as e:
        st.error(f"Failed to generate graph: {e}")


    # agent = create_pandas_dataframe_agent(AzureChatOpenAI(deployment_name="gpt-4", model_name="gpt-4"), df_data)
    # insight_data = agent.run(userInputData,{"df_data": df_data})

def apiDetails():
    #API details - Starts
    TENANT= "XXXXXXXXXXXXXXXXXXXXXXXXXXXX"
    CLIENT_ID = "XXXXXXXXXXXXXXXXXXXXXXXXXXXX"
    CLIENT_SECRET= st.secrets.CLIENT_SECRET
    credential = ClientSecretCredential(TENANT,CLIENT_ID,CLIENT_SECRET)
    VAULT_URL= "https://XXXXXXXXXXXXXXXXXXXXXXXXXXXX.azure.net/"
    client = SecretClient(vault_url=VAULT_URL, credential=credential)
    openai_key = client.get_secret("XXXXXXXXXXXXXXXXXXXXXXXXXXXX")


    os.environ["OPENAI_API_TYPE"] = "azure"
    os.environ["OPENAI_API_BASE"] = "https://XXXXXXXXXXX.openai.azure.com/"
    os.environ["OPENAI_API_KEY"] = openai_key.value
    os.environ["OPENAI_API_VERSION"] = "2023-07-01-preview"
    #Api details - Ends

def downloadText(final_output,finalEmail):
    with open(final_output, 'w') as file:
        file.write(finalEmail)
        with open(final_output, "rb") as file:
            btn = st.download_button(
                label="Download File",
                data=file,
                file_name="Email.txt",
                mime="text/plain"
            )


def main():
   
    #**********************************************************************************#


    file_pathProdDesc = "./temp/file.txt"
    file_pathData = "./temp/file.csv"
    graph_path = "./temp/plot.png"
    final_output = "./temp/Email.txt"


    #**********************************************************************************#


    # streamlit Start


    st.set_page_config(page_title = "E-MAIL MARKETING", page_icon="🚀", layout = "wide")#, initial_sidebar_state = "expanded")
    st.title("GENERATE MARKETING CONTENT")


    uploadProd = ''
    uploadFoot = ''
    prodDesc = 'Empty'
    uploaded_file = ''
    footerInfo = ''
    userInputData = 'No questions from data'
    productDescFile = None
    outputEmail = ''
    global finalEmail
    finalEmail = None
    dataFile = None
    insights = None


    # Email Length
    emailLengthOptions = ['Short','Medium','Long','Enter Word Count']
    emailLength = st.sidebar.selectbox("Email Length",emailLengthOptions,index=1)
    if emailLength == emailLengthOptions[3]:
        emailLength = st.sidebar.text_input("Enter the length of Email")




    # Type of Marketing
    emailMarketingOptions = ["Promotional 🔉", "Newsletter 📰", "Educational 📚","Event 🎉","Not sure ⛔","Please specify if anything else ⌨️",""]
    emailMarketingTyp = st.sidebar.selectbox("Type of marketing content you'd like to create!", emailMarketingOptions[:6], index=0)


    if emailMarketingTyp == emailMarketingOptions[5]:
        emailMarketingTyp = st.sidebar.text_input("Enter the type of marketing content")


    # Product Details
   
    userInputE, userInputD = st.columns(2)
    with userInputE:
        userInputEmail = st.text_area("Outline your email vision",placeholder="Focus on the Problem & Solution...")


    productDescOption = ["Enter Details Manually","Upload a File"]
    productDescBar = st.sidebar.selectbox("What's your preferred way to add description?",productDescOption , index=1)
    if productDescBar == productDescOption[0]:
        prodDesc = st.text_input("Please enter Product Description")
        if prodDesc != 'Empty':
            with open(file_pathProdDesc, 'w') as file:
                file.write(prodDesc)


    elif productDescBar == productDescOption[1]:
        uploadProdtype = st.radio(
        "Upload Product Description",('TEXT','PDF','PPT'), horizontal=True)
        if uploadProdtype == 'PPT':
            productDescFile = st.file_uploader("Upload the file", type=None,key=0)
        elif uploadProdtype == 'TEXT':
            productDescFile = st.file_uploader("Upload the file", type="txt",key=0)
        elif uploadProdtype == 'PDF':
            productDescFile = st.file_uploader("Upload the file", type="pdf",key=0)


    # Dataset


    if "userInputData" not in st.session_state:
        st.session_state.userInputData = False


    uploadData = st.sidebar.radio(
    "Upload the Dataset",('CSV','EXCEL','TEXT'), horizontal=True)
    if uploadData == 'CSV':
        dataFile = st.sidebar.file_uploader("Upload the Dataset", type="csv",key=1)
        with userInputD:
            st.session_state.userInputData = st.text_area("What key information do you need from the CSV?",placeholder="Show the growth of product sales over year...",disabled=not dataFile)
    elif uploadData == 'EXCEL':
        dataFile = st.sidebar.file_uploader("Upload the Dataset", type="xlsx",key=1)
        with userInputD:
            st.session_state.userInputData = st.text_area("What key information do you need from the EXCEL?",placeholder="Show the growth of product sales over year...",disabled=not dataFile)
    elif uploadData == 'TEXT':
        dataFile = st.sidebar.file_uploader("Upload the Dataset", type="txt",key=1)
        with userInputD:
            st.session_state.userInputData = st.text_area("What key information do you need from the TEXT?",placeholder="Show the growth of product sales over year...",disabled=not dataFile)
   
    if productDescFile is not None:
        write_text_file(productDescFile,file_pathProdDesc,uploadProdtype)
   
    if dataFile is not None:
        write_text_file(dataFile,file_pathData,uploadData)


    # Creativity slider
    creativity = st.sidebar.slider(f"How Creative? (0-10)", 0, 10, 3)


    footerOption = ["""The content of this email is confidential and intended for the recipient specified in message only. It is strictly forbidden to share any part of this message with any third party, without a written consent of the sender. If you received this message by mistake, please reply to this message and follow with its deletion, so that we can ensure such a mistake does not occur in the future.""","""It's strictly prohibited to share, copy, print, or otherwise process the content of this email without a written consent from the sender."""]
    legalFooter = st.sidebar.radio("Pick Your Preferred Legal Footer",footerOption)


    if "outputEmail" not in st.session_state:
        st.session_state.outputEmail = False


    if st.button("Time to craft the email!"):
        # Your code to be executed only once goes here
        apiDetails()
        if st.session_state.userInputData:
            insights = dataPrompt(file_pathData,st.session_state.userInputData,graph_path)
            st.image(graph_path, caption='Insight')
            # st.write(dataOutput)
        st.session_state.outputEmail = emailPrompt(file_pathProdDesc,emailMarketingTyp,legalFooter,userInputEmail,emailLength,creativity,insights)
        st.write(st.session_state.outputEmail )
       
    if st.session_state.outputEmail:
        if st.button("Legal Check"):
        # st.button("Legal Check")
            finalEmail = legalPrompt(outputEmail,emailLength)          
            st.write(finalEmail)
            # st.code(finalEmail, language="text")
            downloadText(final_output,finalEmail)


            if st.session_state.userInputData:            
                st.image(graph_path, caption='Insight')
                with open(graph_path, "rb") as file:
                    btn2 = st.download_button(
                        label="Download Image",
                        data=graph_path,
                        file_name="graph.png",  # Customize filename as needed
                        mime="image/png"  # Adjust MIME type based on image format (e.g., JPG, JPEG)
                    )
       
    # streamlit END
    #**********************************************************************************#


if __name__ == '__main__':
    main()





